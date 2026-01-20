from pptx import Presentation
import pdfplumber

MAX_LINES_PER_SLIDE = 10

def pdf_to_ppt(pdf_path, output_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content

    with pdfplumber.open(pdf_path) as pdf:
        for page_no, page in enumerate(pdf.pages):
            text = page.extract_text()

            if not text:
                continue

            lines = [line.strip() for line in text.split("\n") if line.strip()]

            for i in range(0, len(lines), MAX_LINES_PER_SLIDE):
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = f"Page {page_no + 1}"

                content = "\n".join(lines[i:i + MAX_LINES_PER_SLIDE])
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
                text_frame.text = content

    prs.save(output_path)
